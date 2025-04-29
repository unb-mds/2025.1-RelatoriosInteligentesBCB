# Arquivo: criar_apresentacao.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle):
    """Adiciona o slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def add_content_slide(prs, title, bullet_points=None, code=None):
    """Adiciona um slide de conteúdo com título, pontos e opcional código"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    if bullet_points:
        tf = body_shape.text_frame
        tf.clear()
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    
    if code:
        left = Inches(0.5)
        top = Inches(3.0)
        width = Inches(9.0)
        height = Inches(3.0)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        
        p = tf.add_paragraph()
        p.text = code
        font = p.font
        font.size = Pt(10)
        font.name = "Courier New"
        p.line_spacing = 1.0
    
    return slide

def add_image_slide(prs, title, image_path, description=None):
    """Adiciona um slide com título e imagem"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    
    title_shape.text = title
    
    if os.path.exists(image_path):
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(8.0)
        slide.shapes.add_picture(image_path, left, top, width=width)
    
    if description:
        left = Inches(1.0)
        top = Inches(5.0)
        width = Inches(8.0)
        height = Inches(1.0)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = description
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Cria apresentação completa do projeto"""
    prs = Presentation()
    
    # Slide 1: Capa
    add_title_slide(
        prs, 
        "Sistema de Análise de Dados Econômicos com Python e Streamlit",
        "Projeto passo a passo para coleta, visualização e previsão de indicadores econômicos"
    )
    
    # Slide 2: Visão Geral do Projeto
    add_content_slide(
        prs,
        "Visão Geral do Projeto",
        [
            "Objetivo: Criar um sistema para análise de indicadores econômicos brasileiros",
            "Funcionalidades principais:",
            "  • Coleta automática de dados do Banco Central do Brasil",
            "  • Armazenamento em banco de dados SQLite",
            "  • Visualização interativa com Streamlit",
            "  • Previsões usando Machine Learning"
        ]
    )
    
    # Slide 3: Tecnologias Utilizadas
    add_content_slide(
        prs,
        "Tecnologias Utilizadas",
        [
            "Python: Linguagem de programação principal",
            "SQLite: Banco de dados leve e prático",
            "Streamlit: Framework para dashboards interativos",
            "Pandas/NumPy: Manipulação e análise de dados",
            "Plotly/Matplotlib: Visualização de dados",
            "Scikit-learn: Modelos de Machine Learning",
            "Git/GitHub: Controle de versão e colaboração"
        ]
    )
    
    # Slide 4: Estrutura do Projeto
    add_content_slide(
        prs,
        "Estrutura do Projeto",
        [
            "setup.py: Configuração inicial do ambiente",
            "data_collector.py: Coleta de dados das APIs do Banco Central",
            "database_manager.py: Gerenciamento do banco de dados SQLite",
            "app.py: Dashboard principal para visualização",
            "ml_models.py: Implementação dos modelos de Machine Learning",
            "ml_app.py: Dashboard para previsões com ML",
            "main.py: Script principal que integra todos os componentes"
        ]
    )
    
    # Slide 5: Configuração do Ambiente
    add_content_slide(
        prs,
        "Configuração do Ambiente",
        [
            "Passo 1: Instalar Python (versão 3.8+)",
            "Passo 2: Instalar VS Code como ambiente de desenvolvimento",
            "Passo 3: Criar pasta do projeto",
            "Passo 4: Configurar ambiente virtual Python:",
            "  • python -m venv venv",
            "  • source venv/bin/activate (macOS/Linux)",
            "Passo 5: Instalar dependências com setup.py"
        ],
        code="# Trecho de setup.py\ndef install_requirements():\n    packages = [\n        \"pandas\",\n        \"numpy\",\n        \"streamlit\",\n        \"sqlalchemy\",\n        \"scikit-learn\",\n        \"plotly\"\n    ]\n    for package in packages:\n        subprocess.check_call([sys.executable, \"-m\", \"pip\", \"install\", package])"
    )
    
    # Slide 6: Coleta de Dados
    add_content_slide(
        prs,
        "Coleta de Dados (data_collector.py)",
        [
            "Classe BCBDataCollector para acesso às APIs do Banco Central",
            "Indicadores econômicos: IPCA, PIB, Selic, Dívida/PIB, etc.",
            "Processo de requisição e transformação dos dados",
            "Armazenamento temporário em DataFrames do pandas"
        ],
        code="# Trecho de data_collector.py\ndef get_data(self, indicator, start_date=None, end_date=None):\n    \"\"\"Obter dados de um indicador específico\"\"\"\n    serie_id = self.indicators[indicator]\n    url = f\"{self.base_url}.{serie_id}/dados?formato=json\"\n    \n    if start_date:\n        url += f\"&dataInicial={start_date}\"\n    if end_date:\n        url += f\"&dataFinal={end_date}\"\n    \n    response = requests.get(url)\n    data = response.json()"
    )
    
    # Slide 7: Gerenciamento do Banco de Dados
    add_content_slide(
        prs,
        "Gerenciamento do Banco de Dados (database_manager.py)",
        [
            "Classe DatabaseManager para interação com SQLite",
            "Criação automática de tabelas e índices",
            "Métodos para salvar, atualizar e carregar dados",
            "Otimizações para melhor performance"
        ],
        code="# Trecho de database_manager.py\ndef _create_tables(self):\n    \"\"\"Cria as tabelas necessárias no banco de dados\"\"\"\n    conn = sqlite3.connect(self.db_path)\n    cursor = conn.cursor()\n    \n    # Para cada indicador\n    for indicator in indicators:\n        cursor.execute(f'''\n        CREATE TABLE IF NOT EXISTS {indicator} (\n            id INTEGER PRIMARY KEY,\n            date DATE NOT NULL,\n            value FLOAT NOT NULL,\n            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP\n        )''')\n        \n        # Criar índice para a coluna de data\n        cursor.execute(f'CREATE INDEX IF NOT EXISTS idx_{indicator}_date ON {indicator} (date)')"
    )
    
    # Slide 8: Dashboard de Visualização
    add_content_slide(
        prs,
        "Dashboard de Visualização (app.py)",
        [
            "Interface interativa com Streamlit",
            "Seleção de indicadores e período de análise",
            "Gráficos de séries temporais com Plotly",
            "Análise comparativa entre indicadores",
            "Estatísticas descritivas e visualizações"
        ],
        code="# Trecho de app.py\n# Interface principal\nst.title(\"Dashboard Econômico - Dados do Banco Central do Brasil\")\n\n# Sidebar\nst.sidebar.title(\"Opções\")\n\n# Seleção de indicadores\nindicators = st.sidebar.multiselect(\n    \"Selecione os indicadores para visualizar\",\n    list(indicator_names.keys()),\n    default=['ipca', 'selic', 'pib']\n)"
    )
    
    # Slide 9: Modelos de Machine Learning
    add_content_slide(
        prs,
        "Modelos de Machine Learning (ml_models.py)",
        [
            "Classe EconomicPredictor para previsão de indicadores",
            "Preparação de dados para séries temporais",
            "Modelos implementados: Linear, Ridge, Lasso, Random Forest",
            "Avaliação de performance com métricas",
            "Visualização de resultados de previsão"
        ],
        code="# Trecho de ml_models.py\ndef train_model(self, target_indicator, model_type='random_forest', test_size=0.2):\n    \"\"\"Treina um modelo para prever o indicador alvo\"\"\"\n    # Preparar dados\n    X, y, dates = self.prepare_data(target_indicator)\n    \n    # Dividir em treino e teste\n    split_idx = int(len(X) * (1 - test_size))\n    X_train, X_test = X.iloc[:split_idx], X.iloc[split_idx:]\n    y_train, y_test = y.iloc[:split_idx], y.iloc[split_idx:]\n    \n    # Selecionar e treinar o modelo\n    if model_type == 'random_forest':\n        model = RandomForestRegressor(n_estimators=100, random_state=42)"
    )
    
    # Slide 10: Dashboard de Previsões
    add_content_slide(
        prs,
        "Dashboard de Previsões (ml_app.py)",
        [
            "Interface para treinamento de modelos",
            "Seleção de indicador, modelo e período de previsão",
            "Visualização de previsões futuras",
            "Análise de importância de features",
            "Comparação com dados históricos"
        ],
        code="# Trecho de ml_app.py\n# Seleção de indicador\nindicator = st.sidebar.selectbox(\n    \"Selecione o indicador para prever\",\n    list(available_indicators.keys()),\n    format_func=lambda x: available_indicators[x]\n)\n\n# Seleção de modelo\nmodel_type = st.sidebar.selectbox(\n    \"Selecione o modelo\",\n    [\"Linear Regression\", \"Ridge Regression\", \"Lasso Regression\", \"Random Forest\"],\n    index=0\n)"
    )
    
    # Slide 11: Integração
    add_content_slide(
        prs,
        "Integração (main.py)",
        [
            "Script principal que integra todos os componentes",
            "Interface unificada com três opções principais:",
            "  • Coleta de Dados: atualiza o banco de dados",
            "  • Dashboard Econômico: visualiza indicadores",
            "  • Previsões com ML: treina modelos e faz previsões",
            "Facilita a navegação e uso do sistema"
        ],
        code="# Trecho de main.py\n# Menu principal\nst.header(\"Escolha uma opção:\")\n\ncol1, col2, col3 = st.columns(3)\n\nwith col1:\n    st.subheader(\"Coleta de Dados\")\n    if st.button(\"Coletar Dados\"):\n        with st.spinner(\"Coletando dados...\"):\n            collect_data()\n\nwith col2:\n    st.subheader(\"Dashboard Econômico\")\n    if st.button(\"Abrir Dashboard\"):\n        run_dashboard()"
    )
    
    # Slide 12: Controle de Versão e Colaboração
    add_content_slide(
        prs,
        "Controle de Versão e Colaboração",
        [
            "Inicialização do Git para controle de versão",
            "Criação de arquivo .gitignore para ignorar arquivos temporários",
            "Criação de repositório no GitHub para compartilhamento",
            "Instruções para colegas clonarem e executarem o projeto",
            "Fluxo de trabalho colaborativo"
        ],
        code="# Comandos Git básicos\n\n# Inicializar repositório\ngit init\n\n# Adicionar arquivos\ngit add .\n\n# Fazer commit\ngit commit -m \"Versão inicial do projeto\"\n\n# Conectar ao GitHub\ngit remote add origin https://github.com/seu-usuario/projeto-dados-economicos.git\ngit push -u origin main"
    )
    
    # Slide 13: Executando o Projeto
    add_content_slide(
        prs,
        "Executando o Projeto",
        [
            "Passo 1: Ativar o ambiente virtual",
            "  • source venv/bin/activate (macOS)",
            "Passo 2: Verificar/atualizar dependências",
            "  • python setup.py",
            "Passo 3: Coletar dados (se necessário)",
            "  • python data_collector.py",
            "Passo 4: Executar a aplicação principal",
            "  • streamlit run main.py",
            "Comandos alternativos para componentes específicos:"
        ],
        code="# Executar apenas o dashboard de visualização\nstreamlit run app.py\n\n# Executar apenas o dashboard de ML\nstreamlit run ml_app.py\n\n# Para colegas que clonam o repositório\ngit clone https://github.com/seu-usuario/projeto-dados-economicos.git\ncd projeto-dados-economicos\npython -m venv venv\nsource venv/bin/activate  # macOS/Linux\npython setup.py"
    )
    
    # Slide 14: Expansão do Projeto
    add_content_slide(
        prs,
        "Expansão do Projeto",
        [
            "Adição de novos indicadores econômicos",
            "Implementação de modelos mais avançados (LSTM, Prophet)",
            "Criação de alertas para mudanças significativas",
            "Análises mais complexas (correlação, causalidade)",
            "Hospedagem na web (Streamlit Cloud, Heroku)",
            "Automação de coleta periódica de dados",
            "Interface de usuário mais elaborada"
        ]
    )
    
    # Slide 15: Recursos e Referências
    add_content_slide(
        prs,
        "Recursos e Referências",
        [
            "Documentação do Banco Central do Brasil (APIs)",
            "Documentação do Streamlit: https://docs.streamlit.io",
            "Documentação do Scikit-learn: https://scikit-learn.org",
            "Documentação do Pandas: https://pandas.pydata.org/docs/",
            "Documentação do Plotly: https://plotly.com/python/",
            "Tutoriais de Git e GitHub",
            "Livros e cursos sobre análise de dados e séries temporais"
        ]
    )
    
    # Slide 16: Perguntas e Contato
    add_content_slide(
        prs,
        "Perguntas e Contato",
        [
            "Dúvidas frequentes:",
            "  • Como adicionar novos indicadores?",
            "  • Como configurar o ambiente no Windows?",
            "  • Como resolver problemas de coleta de dados?",
            "",
            "Para mais informações, entre em contato:",
            "  • Email: seu-email@exemplo.com",
            "  • GitHub: github.com/seu-usuario",
            "",
            "Obrigado!"
        ]
    )
    
    # Salvar a apresentação
    pptx_path = "Projeto_Dados_Economicos.pptx"
    prs.save(pptx_path)
    print(f"Apresentação criada com sucesso: {pptx_path}")
    
    return pptx_path

if __name__ == "__main__":
    # Verificar se python-pptx está instalado
    try:
        import pptx
    except ImportError:
        print("Instalando biblioteca python-pptx...")
        import subprocess
        import sys
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    
    # Criar a apresentação
    pptx_path = create_presentation()
    
    # Apresentar o caminho completo onde o arquivo foi salvo
    import os
    absolute_path = os.path.abspath(pptx_path)
    print(f"Arquivo salvo em: {absolute_path}")
    
    # Abrir o arquivo (opcional, dependendo do sistema)
    try:
        import platform
        if platform.system() == 'Darwin':  # macOS
            subprocess.call(['open', absolute_path])
        elif platform.system() == 'Windows':
            os.startfile(absolute_path)
        else:  # Linux
            subprocess.call(['xdg-open', absolute_path])
    except:
        print("Arquivo criado, mas não foi possível abri-lo automaticamente.")